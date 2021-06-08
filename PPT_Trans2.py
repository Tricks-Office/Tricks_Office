from google.cloud import translate_v2 as translate
import os
from pptx import Presentation

def GTrans2(file,TPath,RPath,language,Key):
    os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = Key
    trans=translate.Client()
    prs=Presentation(TPath + "/" + file)
    Slides = prs.slides
    for slide in Slides:
        Shapes = slide.shapes
        for shape in Shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    font=paragraph.runs[0].font
                    Fname = font.name
                    Fsize = font.size
                    Fbold = font.bold
                    cur_text = paragraph.text
                    new_text = trans.translate(cur_text, target_language=language)
                    paragraph.text = new_text['translatedText']
                    for run in paragraph.runs:
                        font = run.font
                        font.name = Fname
                        font.size = Fsize
                        font.bold = Fbold
    prs.save(RPath + "/" + file)